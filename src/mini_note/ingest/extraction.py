"""内容提取 — 按文件类型解析文本、生成 ExtractionResult。"""

from dataclasses import dataclass, field
from datetime import datetime, timezone, timedelta
from pathlib import Path

CST = timezone(timedelta(hours=8))


@dataclass
class ExtractionResult:
    """解析结果。

    Attributes:
        source_id: 关联的 source
        status: full/partial/metadata_only/failed
        content: 提取的 Markdown 文本
        extractor: 提取器名称
        extractor_version: 提取器版本
        coverage: 解析覆盖信息
        warnings: 警告列表
        created_at: 时间戳
    """

    source_id: str = ""
    status: str = "full"
    content: str = ""
    extractor: str = ""
    extractor_version: str = ""
    coverage: dict = field(default_factory=dict)
    warnings: list[str] = field(default_factory=list)
    created_at: str = ""

    def __post_init__(self):
        if not self.created_at:
            self.created_at = datetime.now(CST).isoformat()


def extract_markdown(path: Path) -> ExtractionResult:
    """提取 Markdown 文件全文。"""
    content = path.read_text(encoding="utf-8", errors="replace")
    return ExtractionResult(
        content=content,
        extractor="markdown-reader",
        extractor_version="1.0",
        status="full",
        coverage={"bytes_read": path.stat().st_size},
    )


def extract_text(path: Path, max_bytes: int | None = None) -> ExtractionResult:
    """提取纯文本文件内容。支持 UTF-8 BOM。"""
    raw = path.read_bytes()

    # 处理 BOM
    if raw.startswith(b"\xef\xbb\xbf"):
        raw = raw[3:]

    # 超限处理
    if max_bytes and len(raw) > max_bytes:
        raw = raw[:max_bytes]
        status = "partial"
        reason = f"文件超过 {max_bytes} 字节上限，仅提取前 {max_bytes} 字节"
    else:
        status = "full"
        reason = None

    content = raw.decode("utf-8", errors="replace")

    return ExtractionResult(
        content=content,
        extractor="text-reader",
        extractor_version="1.0",
        status=status,
        coverage={
            "bytes_total": path.stat().st_size,
            "bytes_read": len(raw),
            "reason": reason,
        },
    )


def extract_docx(path: Path) -> ExtractionResult:
    """提取 DOCX 文件段落、标题、表格文本。"""
    try:
        from docx import Document
    except ImportError:
        return ExtractionResult(
            status="failed",
            extractor="python-docx",
            extractor_version="n/a",
            warnings=["python-docx 未安装"],
        )

    try:
        doc = Document(str(path))
    except Exception:
        return ExtractionResult(
            status="failed",
            extractor="python-docx",
            extractor_version="n/a",
            warnings=["文件无法解析为 DOCX"],
        )

    parts: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if para.style.name.startswith("Heading"):
            level = para.style.name.split()[-1]
            try:
                lv = int(level)
                parts.append(f"{'#' * lv} {text}")
            except ValueError:
                parts.append(f"**{text}**")
        else:
            parts.append(text)

    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            parts.append(" | ".join(cells))

    return ExtractionResult(
        content="\n\n".join(parts),
        extractor="python-docx",
        extractor_version="1.x",
        status="full",
    )


def extract_pdf(path: Path) -> ExtractionResult:
    """提取 PDF 文件文本（pdfplumber）。"""
    try:
        import pdfplumber
    except ImportError:
        return ExtractionResult(
            status="failed",
            extractor="pdfplumber",
            extractor_version="n/a",
            warnings=["pdfplumber 未安装"],
        )

    try:
        with pdfplumber.open(path) as pdf:
            texts = []
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    texts.append(t)
            content = "\n\n".join(texts)
            status = "metadata_only" if not content.strip() else "full"
        return ExtractionResult(
            content=content,
            extractor="pdfplumber",
            extractor_version="0.11.x",
            status=status,
        )
    except Exception as e:
        return ExtractionResult(
            status="failed",
            extractor="pdfplumber",
            extractor_version="0.11.x",
            warnings=[str(e)],
        )


def extract_xlsx(path: Path) -> ExtractionResult:
    """提取 XLSX 文件所有 sheet 的表格数据为 Markdown 表格。"""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return ExtractionResult(
            status="failed",
            extractor="openpyxl",
            extractor_version="n/a",
            warnings=["openpyxl 未安装"],
        )

    try:
        wb = load_workbook(path, read_only=True, data_only=True)
    except Exception:
        return ExtractionResult(
            status="failed",
            extractor="openpyxl",
            extractor_version="n/a",
            warnings=["文件无法解析为 XLSX"],
        )

    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"## {sheet_name}")
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            parts.append("（空表）")
            continue

        # 第一行作为表头
        header = [str(c) if c is not None else "" for c in rows[0]]
        parts.append("| " + " | ".join(header) + " |")
        parts.append("| " + " | ".join("---" for _ in header) + " |")

        for row in rows[1:]:
            cells = [str(c) if c is not None else "" for c in row]
            # 跳过全空行
            if any(c.strip() for c in cells):
                parts.append("| " + " | ".join(cells) + " |")

    wb.close()
    return ExtractionResult(
        content="\n\n".join(parts),
        extractor="openpyxl",
        extractor_version="3.x",
        status="full",
        coverage={"sheets": len(wb.sheetnames)},
    )


def extract_pptx(path: Path) -> ExtractionResult:
    """提取 PPTX 文件所有幻灯片的文本和表格。"""
    try:
        from pptx import Presentation
    except ImportError:
        return ExtractionResult(
            status="failed",
            extractor="python-pptx",
            extractor_version="n/a",
            warnings=["python-pptx 未安装"],
        )

    try:
        prs = Presentation(str(path))
    except Exception:
        return ExtractionResult(
            status="failed",
            extractor="python-pptx",
            extractor_version="n/a",
            warnings=["文件无法解析为 PPTX"],
        )

    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"## 幻灯片 {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    parts.append(" | ".join(cells))

    return ExtractionResult(
        content="\n\n".join(parts),
        extractor="python-pptx",
        extractor_version="0.6.x",
        status="full",
        coverage={"slides": len(prs.slides)},
    )


def extract_image(path: Path) -> ExtractionResult:
    """提取图片元数据（尺寸、格式、模式）。"""
    try:
        from PIL import Image
    except ImportError:
        return ExtractionResult(
            status="failed",
            extractor="pillow",
            extractor_version="n/a",
            warnings=["Pillow 未安装"],
        )

    try:
        img = Image.open(path)
        width, height = img.size
        fmt = img.format or path.suffix.upper().lstrip(".")
        mode = img.mode
        img.close()

        content = f"""## 图片元数据

| 属性 | 值 |
|------|-----|
| 宽度 | {width} px |
| 高度 | {height} px |
| 格式 | {fmt} |
| 色彩模式 | {mode} |
| 文件大小 | {path.stat().st_size} bytes |
"""
        return ExtractionResult(
            content=content,
            extractor="pillow",
            extractor_version="10.x",
            status="metadata_only",
            coverage={
                "width": width,
                "height": height,
                "format": fmt,
                "mode": mode,
                "size_bytes": path.stat().st_size,
            },
        )
    except Exception as e:
        return ExtractionResult(
            status="failed",
            extractor="pillow",
            extractor_version="10.x",
            warnings=[str(e)],
        )


def extract_code(path: Path) -> ExtractionResult:
    """提取代码文件的结构化摘要（函数、类、导入、行数）。"""
    ext = path.suffix.lower()
    text = path.read_text(encoding="utf-8", errors="replace")
    lines = text.split("\n")
    line_count = len(lines)

    if ext == ".py":
        summary = _parse_python(text)
    else:
        summary = _parse_generic_code(text, ext)

    content = f"""## 代码摘要

**文件**: {path.name}
**语言**: {ext.lstrip('.')}
**总行数**: {line_count}

{summary}
"""
    return ExtractionResult(
        content=content,
        extractor="code-parser",
        extractor_version="1.0",
        status="full",
        coverage={
            "language": ext.lstrip("."),
            "lines": line_count,
        },
    )


def _parse_python(text: str) -> str:
    """使用 AST 解析 Python 代码结构。"""
    import ast

    parts: list[str] = []
    try:
        tree = ast.parse(text)
        for node in ast.iter_child_nodes(tree):
            if isinstance(node, ast.Import):
                names = [a.name for a in node.names]
                parts.append(f"- **导入**: `import {', '.join(names)}`")
            elif isinstance(node, ast.ImportFrom):
                names = [a.name for a in node.names]
                parts.append(f"- **导入**: `from {node.module} import {', '.join(names)}`")
            elif isinstance(node, ast.FunctionDef):
                args = [a.arg for a in node.args.args]
                doc = ast.get_docstring(node)
                parts.append(f"- **函数**: `def {node.name}({', '.join(args)})`")
                if doc:
                    parts.append(f"  - 说明: {doc.split(chr(10))[0][:80]}")
            elif isinstance(node, ast.AsyncFunctionDef):
                args = [a.arg for a in node.args.args]
                parts.append(f"- **异步函数**: `async def {node.name}({', '.join(args)})`")
            elif isinstance(node, ast.ClassDef):
                bases = [ast.unparse(b) for b in node.bases] if node.bases else []
                base_str = f"({', '.join(bases)})" if bases else ""
                parts.append(f"- **类**: `class {node.name}{base_str}`")
                doc = ast.get_docstring(node)
                if doc:
                    parts.append(f"  - 说明: {doc.split(chr(10))[0][:80]}")
                for body_node in node.body:
                    if isinstance(body_node, ast.FunctionDef):
                        a = [a.arg for a in body_node.args.args]
                        parts.append(f"  - **方法**: `def {body_node.name}({', '.join(a)})`")
    except SyntaxError:
        parts.append("- （代码包含语法错误，跳过 AST 解析）")

    return "\n".join(parts) if parts else "（无结构化元素）"


def _parse_generic_code(text: str, ext: str) -> str:
    """用正则提取通用代码结构。"""
    import re

    parts: list[str] = []
    # 查找常见函数定义模式
    func_patterns = [
        r"^\s*(?:def|function|func|fn)\s+(\w+)",  # Python, JS, TS, Go, Rust
        r"^\s*(?:public|private|protected)?\s*(?:static)?\s*\w+\s+(\w+)\s*\([^)]*\)\s*\{",  # Java/C#
    ]
    for pattern in func_patterns:
        for m in re.finditer(pattern, text, re.MULTILINE):
            name = m.group(1)
            if not any(f"`{name}`" in p for p in parts):
                parts.append(f"- **函数/方法**: `{name}()`")

    # 查找类定义
    class_matches = re.findall(r"^\s*(?:class|interface|struct|enum)\s+(\w+)", text, re.MULTILINE)
    for name in class_matches:
        parts.append(f"- **类/结构**: `{name}`")

    # 查找导入
    import_matches = re.findall(r"^\s*(?:import|require|use|from)\s+(.+)", text, re.MULTILINE)
    for imp in import_matches[:10]:
        parts.append(f"- **导入**: `{imp.strip()}`")

    return "\n".join(parts) if parts else "（无识别到的代码结构）"


def extract_media(path: Path) -> ExtractionResult:
    """提取音视频文件的元数据。

    优先尝试 mutagen（纯 Python），如失败则回退到文件基本信息。
    """
    ext = path.suffix.lower()
    size = path.stat().st_size

    # 尝试用 mutagen 提取标签
    metadata = _try_mutagen(path)
    if metadata:
        content = "## 音视频元数据\n\n| 属性 | 值 |\n|------|-----|\n"
        for k, v in metadata.items():
            content += f"| {k} | {v} |\n"
        return ExtractionResult(
            content=content,
            extractor="mutagen",
            extractor_version="1.x",
            status="metadata_only",
            coverage=metadata,
        )

    # 回退：仅文件基本信息
    audio_exts = {".mp3", ".wav", ".flac", ".ogg", ".m4a", ".aac", ".wma"}
    video_exts = {".mp4", ".avi", ".mov", ".mkv", ".webm", ".flv"}
    kind = "音频" if ext in audio_exts else "视频" if ext in video_exts else "未知"

    content = f"""## 媒体文件基本信息

| 属性 | 值 |
|------|-----|
| 文件名 | {path.name} |
| 类型 | {kind} |
| 大小 | {size} bytes ({size / 1024:.1f} KB) |
| 扩展名 | {ext} |
"""
    return ExtractionResult(
        content=content,
        extractor="file-info",
        extractor_version="1.0",
        status="metadata_only",
        coverage={"kind": kind, "size_bytes": size, "extension": ext},
    )


def _try_mutagen(path: Path) -> dict | None:
    """尝试用 mutagen 提取音视频标签。"""
    try:
        from mutagen import File as MutagenFile
    except ImportError:
        return None

    try:
        mf = MutagenFile(str(path))
        if mf is None:
            return None

        info = {}
        if hasattr(mf, "info") and mf.info:
            for attr in ("length", "bitrate", "sample_rate", "channels", "bits_per_sample"):
                val = getattr(mf.info, attr, None)
                if val is not None:
                    if attr == "length":
                        info["时长"] = f"{val:.1f} 秒"
                    elif attr == "bitrate":
                        info["比特率"] = f"{int(val)} bps"
                    elif attr == "sample_rate":
                        info["采样率"] = f"{int(val)} Hz"
                    else:
                        info[attr] = str(val)

        # 标签
        tag_map = {
            "artist": "艺术家",
            "album": "专辑",
            "title": "标题",
            "date": "日期",
            "genre": "流派",
        }
        if hasattr(mf, "tags") and mf.tags:
            for eng, chn in tag_map.items():
                val = mf.tags.get(eng)
                if val:
                    info[chn] = str(val[0]) if isinstance(val, list) else str(val)

        if not info and hasattr(mf, "pprint"):
            info["原始信息"] = mf.pprint()[:500]

        return info if info else None
    except Exception:
        return None


def extract_by_type(path: Path, max_bytes: int | None = None) -> ExtractionResult:
    """根据文件扩展名自动选择提取器。"""
    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {path}")

    ext = path.suffix.lower()

    extractors = {
        ".md": extract_markdown,
        ".txt": lambda p: extract_text(p, max_bytes=max_bytes),
        ".docx": extract_docx,
        ".pdf": extract_pdf,
        ".xlsx": extract_xlsx,
        ".pptx": extract_pptx,
    }

    if ext in extractors:
        return extractors[ext](path)

    image_exts = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".svg"}
    if ext in image_exts:
        return extract_image(path)

    code_exts = {
        ".py", ".js", ".ts", ".go", ".java", ".rs", ".cpp", ".c", ".h",
        ".sh", ".bash", ".yaml", ".yml", ".toml", ".json", ".xml", ".sql",
        ".rb", ".php", ".swift", ".kt", ".scala",
    }
    if ext in code_exts:
        return extract_code(path)

    media_exts = {
        ".mp3", ".wav", ".flac", ".ogg", ".m4a", ".aac",
        ".mp4", ".avi", ".mov", ".mkv", ".webm", ".flv",
    }
    if ext in media_exts:
        return extract_media(path)

    return ExtractionResult(
        status="metadata_only",
        extractor="unknown",
        extractor_version="n/a",
        warnings=[f"不支持的文件类型: {ext}"],
    )
