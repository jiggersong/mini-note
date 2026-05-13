"""
pytest fixtures for mini-note test suite.

每个测试独立运行，通过临时目录隔离 workspace 状态。
"""

import shutil
import tempfile
from pathlib import Path

import pytest


@pytest.fixture
def empty_workspace():
    """创建一个空临时目录，测试后自动清理。"""
    tmp = Path(tempfile.mkdtemp(prefix="mini-note-test-"))
    yield tmp
    shutil.rmtree(tmp, ignore_errors=True)


@pytest.fixture
def tmp_workspace(empty_workspace):
    """创建完整的 mini-note workspace 目录结构。

    包含 meta/、raw/、wiki/、.state/ 目录及默认配置模板。
    """
    ws = empty_workspace

    # meta/
    (ws / "meta").mkdir()
    (ws / "meta" / "purpose.md").write_text("# 测试知识库\n## 目的\n测试用途。\n")
    (ws / "meta" / "config.yaml").write_text(
        "default_scope: shared\nlock_timeout_seconds: 300\n"
    )

    # raw/
    (ws / "raw" / "inbox" / "users").mkdir(parents=True)
    (ws / "raw" / "inbox" / "teams").mkdir(parents=True)
    (ws / "raw" / "archive").mkdir()
    (ws / "raw" / "extracted").mkdir()

    # wiki/
    for sub in ["entities", "concepts", "sources", "synthesis", "queries"]:
        (ws / "wiki" / sub).mkdir(parents=True)
    (ws / "wiki" / "index.md").write_text("# Index\n")
    (ws / "wiki" / "overview.md").write_text("# Overview\n")
    (ws / "wiki" / "log.md").write_text("# Log\n")

    # .state/
    for sub in ["operations", "review_tasks", "health", "staging"]:
        (ws / ".state" / sub).mkdir(parents=True)

    return ws


@pytest.fixture
def sample_md_file(tmp_workspace):
    """在 workspace 的 inbox 中创建一个测试 Markdown 文件。"""
    content = """# 测试文档

## 简介

这是一份测试文档，用于验证 Markdown 摄入功能。

## 核心内容

ECS 突发性能实例适合低负载且偶发突增的场景。

## 注意事项

- 不要在生产环境直接测试
- 务必备份数据
"""
    path = tmp_workspace / "raw" / "inbox" / "users" / "test-doc.md"
    path.write_text(content, encoding="utf-8")
    return path


@pytest.fixture
def sample_txt_file(tmp_workspace):
    """创建一个纯文本测试文件。"""
    content = "纯文本文件内容。\n第二行。\n第三行包含关键数字：42。\n"
    path = tmp_workspace / "raw" / "inbox" / "users" / "test-doc.txt"
    path.write_text(content, encoding="utf-8")
    return path


@pytest.fixture
def sample_docx_file(tmp_workspace):
    """创建一个最小可解析的 DOCX 测试文件。"""
    try:
        from docx import Document
    except ImportError:
        pytest.skip("python-docx 未安装")

    path = tmp_workspace / "raw" / "inbox" / "users" / "test-doc.docx"
    doc = Document()
    doc.add_heading("测试 DOCX 文档", level=1)
    doc.add_paragraph("这是一个段落，包含一些测试文本。")
    doc.add_paragraph("第二个段落，包含关键信息：最大连接数为 1000。")
    doc.save(str(path))
    return path


@pytest.fixture
def sample_pdf_file(tmp_workspace):
    """创建一个最小可解析的 PDF 测试文件。"""
    # pdfplumber 需要真实 PDF，用简单的 PDF 生成方式
    path = tmp_workspace / "raw" / "inbox" / "users" / "test-doc.pdf"
    # 构造最小合法 PDF
    pdf_content = (
        b"%PDF-1.4\n"
        b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
        b"2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n"
        b"3 0 obj<</Type/Page/MediaBox[0 0 612 792]/Parent 2 0 R/Resources<<>>>>endobj\n"
        b"xref\n0 4\n0000000000 65535 f \n0000000009 00000 n \n0000000058 00000 n \n0000000115 00000 n \n"
        b"trailer<</Size 4/Root 1 0 R>>\n"
        b"startxref\n190\n%%EOF\n"
    )
    path.write_bytes(pdf_content)
    return path


@pytest.fixture
def sample_xlsx_file(tmp_workspace):
    """创建一个最小可解析的 XLSX 测试文件。"""
    try:
        from openpyxl import Workbook
    except ImportError:
        pytest.skip("openpyxl 未安装")

    path = tmp_workspace / "raw" / "inbox" / "users" / "test-data.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "ECS 规格"
    ws["A1"] = "实例规格"
    ws["B1"] = "vCPU"
    ws["C1"] = "内存(GB)"
    ws["D1"] = "最大带宽(Mbps)"
    ws["A2"] = "ecs.t6-c1m1"
    ws["B2"] = 1
    ws["C2"] = 1
    ws["D2"] = 0.5
    ws["A3"] = "ecs.g6.large"
    ws["B3"] = 2
    ws["C3"] = 8
    ws["D3"] = 1.5
    wb.save(str(path))
    return path


@pytest.fixture
def sample_pptx_file(tmp_workspace):
    """创建一个最小可解析的 PPTX 测试文件。"""
    try:
        from pptx import Presentation
    except ImportError:
        pytest.skip("python-pptx 未安装")

    path = tmp_workspace / "raw" / "inbox" / "users" / "test-slides.pptx"
    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    if slide1.shapes.title:
        slide1.shapes.title.text = "ECS 性能优化方案"
    # 找一个有 subtitle 的 shape
    for shape in slide1.shapes:
        if shape.has_text_frame and shape.text_frame.text == "":
            shape.text_frame.text = "突发性能实例的适用场景与限制"
            break

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    if slide2.shapes.title:
        slide2.shapes.title.text = "关键要点"
    body = slide2.shapes.add_textbox(100, 150, 500, 300)
    tf = body.text_frame
    tf.text = "最大连接数限制为 1000"
    p = tf.add_paragraph()
    p.text = "适合低负载且偶发突增的场景"

    prs.save(str(path))
    return path


@pytest.fixture
def sample_image_file(tmp_workspace):
    """创建一个最小 PNG 测试图片。"""
    try:
        from PIL import Image
    except ImportError:
        pytest.skip("Pillow 未安装")

    path = tmp_workspace / "raw" / "inbox" / "users" / "test-image.png"
    img = Image.new("RGB", (640, 480), color=(73, 109, 137))
    img.save(str(path), format="PNG")
    return path


@pytest.fixture
def sample_python_file(tmp_workspace):
    """创建一个 Python 测试代码文件。"""
    code = '''"""
ECS 性能分析模块 — 计算突发性能实例的基线积分。
"""

import math
from dataclasses import dataclass
from typing import Optional


@dataclass
class InstanceSpec:
    """实例规格描述。"""
    instance_type: str
    vcpu: int
    memory_gb: float
    baseline_gops: float
    credit_accumulation: float = 1.0

    def compute_credit_balance(self, duration_seconds: int) -> float:
        """计算积分余额。

        Args:
            duration_seconds: 持续时间（秒）

        Returns:
            积分余额
        """
        if duration_seconds <= 0:
            raise ValueError("duration_seconds 必须为正数")
        return self.baseline_gops * duration_seconds * self.credit_accumulation


def analyze_instances(specs: list[InstanceSpec]) -> dict:
    """分析实例列表的总体特征。"""
    total_vcpu = sum(s.vcpu for s in specs)
    total_memory = sum(s.memory_gb for s in specs)
    return {
        "instance_count": len(specs),
        "total_vcpu": total_vcpu,
        "total_memory_gb": total_memory,
        "avg_credit": total_vcpu / len(specs) if specs else 0,
    }
'''
    path = tmp_workspace / "raw" / "inbox" / "users" / "ecs_analyzer.py"
    path.write_text(code, encoding="utf-8")
    return path


@pytest.fixture
def sample_mp3_file(tmp_workspace):
    """创建一个最小合法 MP3 文件（包含 ID3 头和有效帧）。"""
    path = tmp_workspace / "raw" / "inbox" / "users" / "test-audio.mp3"
    # 最小 MP3: ID3v2 头 + 同步字 + 有效帧头
    # 构造一个最简合法 MP3 帧 (MPEG1 Layer3 128kbps 44100Hz)
    import struct
    data = bytearray()
    # ID3v2 头（最小）
    data += b"ID3\x03\x00\x00\x00\x00\x0a\x00" + b"\x00" * 10
    # MPEG 帧头: 0xFFFB9000 (MPEG1 Layer3 128kbps 44100Hz stereo)
    frame_header = struct.pack(">I", 0xFFFB9000)
    data += frame_header
    # 帧体 (417 bytes for 128kbps/44100)
    data += b"\x00" * 417
    path.write_bytes(bytes(data))
    return path


@pytest.fixture
def sample_large_text_file(tmp_workspace):
    """创建一个超过 2MB 上限的大文本文件。"""
    path = tmp_workspace / "raw" / "inbox" / "users" / "large-doc.txt"
    # 写入超过 2MB 的内容
    chunk = "大数据量测试行。\n" * 10000
    with path.open("w", encoding="utf-8") as f:
        for _ in range(50):
            f.write(chunk)
    return path


@pytest.fixture
def sample_source_yaml():
    """返回一份合法的 source.yaml 样例数据。"""
    return {
        "source_id": "src-20260513-120000-a1b2",
        "original_name": "ecs-guide.pdf",
        "stored_path": "raw/archive/src-20260513-120000-a1b2/original.pdf",
        "sha256": "e3b0c44298fc1c149afbf4c8996fb92427ae41e4649b934ca495991b7852b855",
        "owner_id": "user-default",
        "scope": "shared",
        "media_type": "application/pdf",
        "size_bytes": 1234567,
        "ingestion_status": "full",
        "created_at": "2026-05-13T12:00:00+08:00",
    }


@pytest.fixture
def sample_claim_data():
    """返回一份合法的 claim 样例数据。"""
    return {
        "claim_id": "claim-20260513-120100-0001",
        "source_id": "src-20260513-120000-a1b2",
        "text": "ECS 突发性能实例适合低负载且偶发突增的场景。",
        "locator": "page=6 paragraph=3",
        "quote_hash": "sha256:d34db33f",
        "extraction_method": "pdf_text",
        "confidence": 0.88,
        "status": "active",
        "verified_at": "2026-05-13T12:01:00+08:00",
    }


@pytest.fixture
def sample_operation_manifest():
    """返回一份合法的 operation manifest 样例数据。"""
    return {
        "operation_id": "op-20260513-120000-a1b2",
        "type": "ingest",
        "status": "planned",
        "source_ids": ["src-20260513-120000-a1b2"],
        "planned_changes": [
            {"action": "create_page", "path": "wiki/sources/src-20260513-120000-a1b2.md"},
            {"action": "update_page", "path": "wiki/index.md"},
        ],
        "validation": {
            "frontmatter_valid": True,
            "wikilinks_valid": True,
            "claims_have_sources": True,
        },
    }


# ============================================================
# OSS 可用性检测
# ============================================================

@pytest.fixture(scope="session")
def oss_configured() -> bool:
    """检测当前环境是否配置了 OSS 凭证。

    OSS 测试需要真实凭证，无凭证时跳过。
    """
    import os
    required = ["OSS_ENDPOINT", "OSS_BUCKET", "OSS_ACCESS_KEY_ID", "OSS_ACCESS_KEY_SECRET"]
    return all(os.getenv(k) for k in required)


requires_oss = pytest.mark.skipif(
    not all(
        __import__("os").getenv(k)
        for k in ["OSS_ENDPOINT", "OSS_BUCKET", "OSS_ACCESS_KEY_ID", "OSS_ACCESS_KEY_SECRET"]
    ),
    reason="OSS 凭证未配置，跳过云端测试",
)
